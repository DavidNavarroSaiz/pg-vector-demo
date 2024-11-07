import openai
import os
import tempfile
import base64
from dotenv import load_dotenv
from moviepy.editor import VideoFileClip
from youtube_transcript_api import YouTubeTranscriptApi
from langchain_community.document_loaders import PyMuPDFLoader, Docx2txtLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter

from pptx import Presentation
from pathlib import Path

import re
from PIL import Image

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

class UniversalDocumentProcessor:
    def __init__(self):
        self.input_token_cost = 0.150 / 1_000_000
        self.output_token_cost = 0.600 / 1_000_000
        self.transcription_cost_per_minute = 0.006
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=700,
            chunk_overlap=50,
        )

    def load_document(self, doc_path: str):
        
        if self.is_url(doc_path):
            if "youtube.com" in doc_path or "youtu.be" in doc_path:
                return self.process_youtube_video(doc_path)
            else:
                raise ValueError("URL provided is not a supported format.")
            
        doc_path = Path(doc_path).resolve()
        
        if not Path(doc_path).is_file():
            raise ValueError(f"File path {doc_path} is not a valid file.")
        doc_path = str(doc_path)
        if doc_path.endswith(".pdf"):
            loader = PyMuPDFLoader(doc_path)
        elif doc_path.endswith('.docx') or doc_path.endswith('.doc'):
            loader = Docx2txtLoader(doc_path)
        elif doc_path.endswith('.txt'):
            loader = TextLoader(doc_path)
        elif doc_path.endswith('.pptx'):
            return self.load_pptx(doc_path)
        elif doc_path.endswith('.mp4'):
            return self.process_video(doc_path)
        elif doc_path.endswith(('.jpg', '.jpeg', '.png')):
            return self.process_image(doc_path)
        
        else:
            raise ValueError("Unsupported file format")
        
        docs = loader.load()
        return "\n\n".join(doc.page_content for doc in docs)

    def is_url(self, path):
        # Check if path is a URL
        url_regex = re.compile(r'^(https?://)?(www\.)?([a-zA-Z0-9_-]+)+(\.[a-zA-Z]+)+(/[\w#!:.?+=&%@!\-]*)?$')
        return re.match(url_regex, path) is not None

    def load_pptx(self, path):
        prs = Presentation(path)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return "\n\n".join(content)

    def encode_image(self, image_path):
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')

    def process_image(self, image_path):
        # Get the resolution of the image
        with Image.open(image_path) as img:
            width, height = img.size

        # Encode image to base64 for GPT processing
        base64_image = self.encode_image(image_path)
        prompt = self.get_image_prompt()
        result, prompt_tokens, completion_tokens = self.ask_gpt(prompt, image=base64_image)
        model_cost = self.calculate_model_cost(prompt_tokens, completion_tokens)
        
        return {
            "original_text": "Image content processed",
            "summary": result,
            "resolution": f"{width}x{height} pixels",
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "cost": model_cost
        }

    def process_video(self, video_path):
        audio_path = self.extract_audio_from_video(video_path)
        transcript_text = self.transcribe_audio_whisper(audio_path)
        summary, prompt_tokens, completion_tokens = self.summarize_text(transcript_text)
        audio_duration_minutes = self.get_audio_duration_in_minutes(video_path)
        transcription_cost = self.calculate_whisper_cost(audio_duration_minutes)
        model_cost = self.calculate_model_cost(prompt_tokens, completion_tokens)
        
        total_cost = transcription_cost + model_cost
        return {
            "original_text": transcript_text,
            "summary": summary,
            "audio_duration_minutes": audio_duration_minutes,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "cost": total_cost
        }

    def process_youtube_video(self, video_url):
        video_id = video_url.split('v=')[-1].split('&')[0]
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = " ".join([entry['text'] for entry in transcript])
        summary, prompt_tokens, completion_tokens = self.summarize_text(transcript_text)
        model_cost = self.calculate_model_cost(prompt_tokens, completion_tokens)

        return {
            "original_text": transcript_text,
            "summary": summary,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "cost": model_cost
        }

    def extract_audio_from_video(self, video_path):
        video = VideoFileClip(video_path)
        temp_audio_path = tempfile.mktemp(suffix=".wav")
        video.audio.write_audiofile(temp_audio_path)
        return temp_audio_path

    def transcribe_audio_whisper(self, audio_path):
        with open(audio_path, 'rb') as audio_file:
            response = openai.audio.transcriptions.create(model = "whisper-1", file=audio_file)
        return response.text

    def get_audio_duration_in_minutes(self, video_path):
        video = VideoFileClip(video_path)
        return video.duration / 60

    def calculate_whisper_cost(self, audio_duration_minutes):
        return audio_duration_minutes * self.transcription_cost_per_minute

    def ask_gpt(self, text, image=None):
       # Construct the message structure with specified format
        message_content = [
            {"type": "text", "text": text}
        ]
        
        if image:
            message_content.append(
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{image}",
                        "detail": "low"
                    }
                    
                }
            )

        # Prepare the messages for the OpenAI API request
        messages = [
            {
                "role": "user",
                "content": message_content
            }
        ]
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages
        )
        answer = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens
        return answer, prompt_tokens, completion_tokens

    def calculate_model_cost(self, input_tokens, output_tokens):
        input_cost = input_tokens * self.input_token_cost
        output_cost = output_tokens * self.output_token_cost
        return input_cost + output_cost

    def summarize_text(self, text):
        prompt = self.get_text_prompt(text)
        summary, prompt_tokens, completion_tokens = self.ask_gpt(prompt)
        return summary, prompt_tokens, completion_tokens

    def get_text_prompt(self, text):
        return f"""
        You are an AI assistant specialized in analyzing resources for recommendation systems. Given the text provided, your task is to produce a summary that is highly relevant for recommending the resource in specific contexts. Please perform the following:

        1. **Key Points and Description**: Identify and describe the core themes and essential information conveyed in the resource. Focus on aspects that are uniquely valuable or distinctive.

        2. **Summary for Recommendation**: Provide a but contextually rich, summary that captures the document’s relevance and potential applications for various audiences.

        3. **Important Components and Highlights**: List the key components, highlights, and distinctive elements that make this document useful. Include specifics on any data, insights, or unique perspectives.

        4. **Ideal Contexts for Use**: Suggest precise scenarios or user needs where this document would be particularly beneficial. Focus on practical applications and how it meets certain informational gaps or needs.

        5. **Detailed Recommendation Scenarios**: Describe in detail the scenarios, industries, or roles for which this document would be highly recommendable. Explain why it would be valuable in each case.

        Text to analyze:
        {text}
        """

    def get_image_prompt(self):
        return """
        You are an AI assistant specialized in analyzing images with a focus on their applicability and relevance for recommendation systems. Given the image provided, perform the following analyses to maximize the utility of the image in specific contexts:

        1. **Detailed Description**: Describe the image comprehensively, highlighting all visible elements and unique characteristics that may be relevant in a recommendation scenario.

        2. **Text Transcription (if applicable)**: If the image includes any text, transcribe it precisely, ensuring the text is captured in the correct structure and order.

        3. **Identification of Key Elements**: Pinpoint the essential components, unique features, and any elements that may serve as focal points in different contexts.

        4. **Contextual Relevance**: Outline specific contexts or use cases where this image would be particularly beneficial. Think of industry, roles, and specific situations.

        5. **Value Proposition**: Explain the utility and added value of the image, considering why it would be relevant for certain users or industries.

        6. **Ideal Target Audience**: Identify the ideal audience for this image, specifying relevant sectors, roles, or interests. This may include scenarios where it serves educational, promotional, or informational purposes.
       
        7.  **Detailed Recommendation Scenarios**: Describe in detail the scenarios, industries, or roles for which this image would be highly recommendable. Explain why it would be valuable in each case.


        """

    def process(self, doc_path=None):
        if doc_path:
            original_text = self.load_document(doc_path)
            if isinstance(original_text, dict):  # Already processed video or image
                return original_text

        else:
            raise ValueError("Invalid input. Provide a document path of a image, video, or YouTube URL to process.")

        summary, prompt_tokens, completion_tokens = self.summarize_text(original_text)
        model_cost = self.calculate_model_cost(prompt_tokens, completion_tokens)

        return {
            "original_text": original_text,
            "summary": summary,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "cost": model_cost
        }
    def get_embedding(self,text):
        response = openai.embeddings.create(input=text, model="text-embedding-3-small")
        return response.data[0].embedding
    
    def split_docs(self, docs):
        chunks = []
        if docs:
            chunks = self.text_splitter.split_documents(docs)
        return chunks
if __name__ == "__main__":
    processor = UniversalDocumentProcessor()
    doc_path = "./iceberg english-01.jpg"  
    result = processor.process(doc_path)
    
    # Organized printing of the results with checks for optional keys
    print("\n" + "=" * 50)
    print("Original Text")
    print("=" * 50)
    print(result["original_text"][:1000] + "...\n" if len(result["original_text"]) > 1000 else result["original_text"])

    print("\n" + "=" * 50)
    print("Summary")
    print("=" * 50)
    print(result["summary"])

    # Check and print optional details like resolution and audio duration
    if "resolution" in result:
        print("\n" + "=" * 50)
        print("Image Details")
        print("=" * 50)
        print(f"Resolution: {result['resolution']}")
        
    if "audio_duration_minutes" in result:
        print("\n" + "=" * 50)
        print("Audio Details")
        print("=" * 50)
        print(f"Audio Duration: {result['audio_duration_minutes']} minutes")

    print("\n" + "=" * 50)
    print("Token Details")
    print("=" * 50)
    print(f"Prompt Tokens: {result.get('prompt_tokens', 'N/A')}")
    print(f"Completion Tokens: {result.get('completion_tokens', 'N/A')}")

    print("\n" + "=" * 50)
    print("Cost Details")
    print("=" * 50)
    print(f"GPT Model Cost: ${result['cost']:.6f}")

    print("=" * 50)